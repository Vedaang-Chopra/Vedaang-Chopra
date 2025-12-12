import os
from pptx import Presentation
from datetime import date

# Configuration
SOURCE_DIR = "content/blog/raw_content"
OUTPUT_DIR = "content/blog/markdown"

# Ensure output directory exists
os.makedirs(OUTPUT_DIR, exist_ok=True)

def extract_presentation(filepath):
    prs = Presentation(filepath)
    filename = os.path.basename(filepath)
    
    # Infer title from filename first, then try first slide
    title = os.path.splitext(filename)[0]
    title = title.replace("CS 8803", "").replace("CS8803", "").replace("-", " ").strip()
    
    slug = title.lower().replace(" ", "-").replace("&", "and")
    
    # Relative path for the download link (assuming content/blog/raw_content is in STATIC_PATHS)
    ppt_static_path = f"{{attach}}../raw_content/{filename}"

    md_content = f"""Title: {title}
Date: {date.today()}
Category: Research Paper Analysis
Slug: {slug}
Summary: Detailed analysis and presentation notes for {title}.

<div class="download-box" style="margin-bottom: 2rem; padding: 1rem; background: var(--btn-bg); border-radius: 8px; display: inline-block;">
    <a href="{ppt_static_path}" style="text-decoration: none; font-weight: bold;">
        📥 Download Original Slides (PPTX)
    </a>
</div>

"""
    
    for i, slide in enumerate(prs.slides):
        slide_title = ""
        slide_text = []
        
        # Get Title
        if slide.shapes.title:
            slide_title = slide.shapes.title.text
        
        # Get Text from other shapes
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # Clean up bullet points slightly
                    text = text.lstrip("- ").lstrip("• ")
                    slide_text.append(text)
        
        # Get Notes
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            
        # Format as blog section
        if slide_title:
            md_content += f"## {slide_title}\n\n"
        
        # Content flow: Notes first (often more conversational), then visual bullets as key points
        if notes:
            # Basic cleanup of notes to look like paragraphs
            notes_clean = notes.replace("\n", " ")
            md_content += f"{notes_clean}\n\n"

        if slide_text:
            md_content += "**Key Takeaways:**\n"
            for line in slide_text:
                md_content += f"*   {line}\n"
            md_content += "\n"
            
    return filename, md_content

files = [f for f in os.listdir(SOURCE_DIR) if f.endswith(".pptx")]

for f in files:
    full_path = os.path.join(SOURCE_DIR, f)
    original_name, content = extract_presentation(full_path)
    
    # Clean filename for output
    safe_name = original_name.replace(" ", "_").replace(".pptx", ".md")
    output_path = os.path.join(OUTPUT_DIR, safe_name)
    
    with open(output_path, "w") as out:
        out.write(content)
    
    print(f"Generated {output_path}")
